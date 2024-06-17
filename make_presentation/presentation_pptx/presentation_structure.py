from __future__ import annotations

import logging
import os.path
from typing import TYPE_CHECKING

from pptx import Presentation

from make_presentation.config import ENDING_PRESENTATION_TEXT, path_to_template
from make_presentation.DTO import ImageDTO
from make_presentation.templates.get_settings import (
    get_slides_foreground_pictures_setting, get_slides_pictures_setting)
from make_presentation.templates.template_config import (TEXT_COLOR, TEXT_FONT,
                                                         TEXT_FONT_SETTINGS,
                                                         TEXT_FONT_SIZE)

from .slide import Slide

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationClass

    from make_presentation.DTO import SlideDTO


logger = logging.getLogger(__name__)


class PresentationTemplate:
    """
    Создаёт pptx файл из шаблона указанного в roots.path_to_template.
    name - название призентации
    title - Заголовки призентации
    text - Тексты слайдов призентации
    img - Изображения для слайдов
    """

    def __init__(
        self,
        name: str,
        title: list[str],
        text: list[str],
        images: list[list[ImageDTO]],
        settings: dict[str, dict[str, str]],
        opening_prentation_theme_title: bool,
        ending_presentation_status: bool
    ) -> None:
        self.presentation_name = name
        self.slides_title = title
        self.slides_text = text
        self.slides_count = len(title)
        self.images = images

        self.settings = settings["PRESENTATION_SETTING"]

        if opening_prentation_theme_title:
            self.slides_title = list([name] + self.slides_title)
            self.slides_text = list([""] + self.slides_text)
            self.images = [[ImageDTO(image=None, path="", description="")]] + self.images
            self.slides_count += 1

        if ending_presentation_status:
            finish_text: str = ENDING_PRESENTATION_TEXT
            self.slides_title = list(self.slides_title + [finish_text])
            self.slides_count += 1

        self.pictures_setting = get_slides_pictures_setting(
           template_name=self.settings["TEMPLATE_NAME"],
           num_slides=self.slides_count,
        )

        self.foreground_pictures_setting = get_slides_foreground_pictures_setting(
           template_name=self.settings["TEMPLATE_NAME"],
           num_slides=self.slides_count,
        )

    def __create_presentation_construction(self) -> PresentationClass:
        """
        To create presentation object from existing presentation template.
        """

        search_template = os.path.join(
            path_to_template,
            self.settings["TEMPLATE_NAME"],
            f"_{self.slides_count}.pptx",
        )

        self.presentation: PresentationClass = Presentation(search_template)
        logger.info("A presentation template has been created.")

        return self.presentation

    def create_presentation(self) -> list[SlideDTO]:
        """
        To create presentation based on a particular template.
        """

        self.__create_presentation_construction()

        slides_info: list[SlideDTO] = []

        for slide_number in range(len(self.presentation.slides)):
            title = ""
            if len(self.slides_title) > slide_number:
                title = self.slides_title[slide_number]

            text = ""
            if len(self.slides_text) > slide_number:
                text = self.slides_text[slide_number]

            img = None
            if len(self.images) > slide_number:
                img = self.images[slide_number]

            if slide_number == 0:
                slide_type = "INITIAL"
            elif slide_number == len(self.presentation.slides) - 1:
                slide_type = "END"
            else:
                slide_type = "USUAL"

            slide = Slide(
                slide=self.presentation.slides[slide_number],
                title=title,
                text=text,
                img=img,
                slide_type=slide_type,
                text_font=TEXT_FONT.get(self.settings["TEMPLATE_NAME"]),
                text_font_settings=TEXT_FONT_SETTINGS.get(self.settings["TEMPLATE_NAME"]),
                text_font_size=TEXT_FONT_SIZE.get(self.settings["TEMPLATE_NAME"]),
                text_color=TEXT_COLOR.get(self.settings["TEMPLATE_NAME"]),
                pictures_setting=(
                    self.pictures_setting[slide_number - 1]
                    if slide_type == "USUAL" else None
                ),
                foreground_pictures_setting=(
                    self.foreground_pictures_setting[slide_number - 1]
                    if slide_type == "USUAL" else None
                )
            )

            slide_dto = slide.make_slide()
            slides_info.append(slide_dto)

        return slides_info

    def save_presentation(self, file_save_path: str) -> None:
        self.presentation.save(file_save_path)
        logger.info(f"A presentation has been saved into {file_save_path}.")
