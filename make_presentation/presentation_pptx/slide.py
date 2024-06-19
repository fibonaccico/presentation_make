from __future__ import annotations

import logging
import os
from io import BytesIO
from typing import TYPE_CHECKING, Any

from pptx.dml.color import RGBColor

from make_presentation.config import (DEFAULT_TEXT_FONT,
                                      DEFAULT_TEXT_FONT_SETTINGS,
                                      DEFAULT_TEXT_SIZE,
                                      path_to_foreground_image)
from make_presentation.DTO import ImageInfoDTO, SlideDTO
from make_presentation.errors import PicturesNumberError

from .image_corrector import ImageCorrector

if TYPE_CHECKING:
    from PIL.Image import Image
    from pptx.slide import Slide as PptxSlide

    from make_presentation.DTO import ImageDTO


logger = logging.getLogger(__name__)


class Slide:
    def __init__(
        self,
        slide: PptxSlide,
        title: str,
        text: str,
        img: list[ImageDTO] | None,
        slide_type: str,
        text_font: dict[str, dict[str, str]] | None,
        text_font_settings: dict[str, dict[str, dict[str, bool]]] | None,
        text_font_size: dict[str, dict[str, int]] | None,
        text_color: dict[str, dict[str, list[int]]] | None,
        pictures_setting: list[dict[str, str]] | None,
        foreground_pictures_setting: list[str] | None
    ) -> None:
        self.title = title
        self.text = text
        self.img = img
        self.slide = slide
        self.slide_type = slide_type
        self.text_font = text_font
        self.text_font_settings = text_font_settings
        self.text_font_size = text_font_size
        self.text_color = text_color
        self.pictures_setting = pictures_setting
        self.foreground_pictures_setting = foreground_pictures_setting

    def make_slide(self) -> SlideDTO:
        num_pic = 0

        for shape in self.slide.shapes:
            if shape.has_text_frame:
                if shape.text == "TITLE" and self.title is not None:
                    self.__add_text_to_placeholder(
                        text_placeholder=shape,
                        text=self.title,
                        text_font=(
                            self.text_font[self.slide_type]["TITLE"]
                            if self.text_font else DEFAULT_TEXT_FONT
                        ),
                        text_font_settings=(
                            self.text_font_settings[self.slide_type]["TITLE"]
                            if self.text_font_settings else DEFAULT_TEXT_FONT_SETTINGS
                        ),
                        text_font_size=(
                            self.text_font_size[self.slide_type]["TITLE"]
                            if self.text_font_size else DEFAULT_TEXT_SIZE
                        ),
                        text_color=(
                            self.text_color[self.slide_type]["TITLE"]
                            if self.text_color else None
                        )
                    )
                elif shape.text == "TEXT" and self.text is not None:
                    self.__add_text_to_placeholder(
                        text_placeholder=shape,
                        text=self.text,
                        text_font=(
                            self.text_font[self.slide_type]["TEXT"]
                            if self.text_font else DEFAULT_TEXT_FONT
                        ),
                        text_font_settings=(
                            self.text_font_settings[self.slide_type]["TEXT"]
                            if self.text_font_settings else DEFAULT_TEXT_FONT_SETTINGS
                        ),
                        text_font_size=(
                            self.text_font_size[self.slide_type]["TEXT"]
                            if self.text_font_size else DEFAULT_TEXT_SIZE
                        ),
                        text_color=(
                            self.text_color[self.slide_type]["TEXT"]
                            if self.text_color else None
                        )
                    )
                elif (
                    shape.text == "PIC"
                    and self.img is not None
                    and len(self.img) > num_pic
                    and self.pictures_setting is not None
                ):
                    self.__add_picture(
                        shape=shape,
                        num_pic=num_pic,
                        settings=self.pictures_setting[num_pic],
                    )
                    num_pic += 1

        if self.foreground_pictures_setting is not None:
            self.__add_foreground_images(names=self.foreground_pictures_setting)

        if self.img:
            images: list[ImageInfoDTO] = []
            for image_dto in self.img:
                image_info = ImageInfoDTO(
                    path=image_dto.path, description=image_dto.description
                )
            images.append(image_info)

            return SlideDTO(title=self.title, text=self.text, images=images)

        return SlideDTO(title=self.title, text=self.text, images=None)

    def __add_text_to_placeholder(
        self,
        text_placeholder: Any,
        text: str,
        text_font: str,
        text_font_settings: dict[str, bool],
        text_font_size: int,
        text_color: list[int] | None
    ) -> None:
        text_placeholder.text_frame.clear()

        paragraph = text_placeholder.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        font = run.font
        try:
            text_placeholder.text_frame.fit_text(
                font_family=text_font,
                max_size=text_font_size,
                bold=text_font_settings["BOLD"],
                italic=text_font_settings["ITALIC"],
                font_file=None,
            )

            if text_color:
                font.color.rgb = RGBColor(text_color[0], text_color[1], text_color[2])

        except OSError:
            logger.info("Could not set text fonts because of unsuppored OS.")

    def __add_picture(self, shape: Any, num_pic: int, settings: dict[str, str]) -> None:
        if self.img:
            pic: Image = self.img[num_pic].image

            if len(self.img) < num_pic:
                raise PicturesNumberError("Wrong picteres number.")

        # Удаляет и восстанавливает объект картинки
        shape.element.getparent().remove(shape.element)
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        if pic:
            # Кoрректирует форму картинки для вставки в презентацию
            img_corrector = ImageCorrector(pillow_img=pic, setting=settings)
            pic = img_corrector.correct()

            # Получите байтовые данные изображения
            image_data = BytesIO()
            pic.save(image_data, format="PNG")
            image_data.seek(0)

            # Вставляем в призентацию
            self.slide.shapes.add_picture(image_data, left, top, width, height)

    def __add_foreground_images(self, names: list[str]) -> None:
        for shape in self.slide.shapes:
            if shape.has_text_frame:
                if shape.text in names:
                    # Удаляет и восстанавливает объект картинки
                    shape.element.getparent().remove(shape.element)
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    # Вставляем в призентацию
                    self.slide.shapes.add_picture(
                        os.path.join(path_to_foreground_image, shape.text),
                        left,
                        top,
                        width,
                        height,
                    )
